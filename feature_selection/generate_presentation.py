"""
Generate a PowerPoint presentation from feature_selection reports.

Creates: feature_selection_outputs/feature_selection_presentation.pptx

Usage (project root):
  python feature_selection/generate_presentation.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import textwrap
import sys

ROOT = Path(__file__).resolve().parent.parent
FS_DIR = ROOT / "feature_selection"
OUT_DIR = ROOT / "feature_selection_outputs"
X_DIR = OUT_DIR / "cross_appliance"
PPT_PATH = X_DIR / "feature_selection_presentation.pptx"

FIGS = [
    ("fig01_pipeline_counts.png", "Pipeline outcome per appliance"),
    ("fig02_stability_heatmap.png", "Feature stability heatmap (50×5)"),
    ("fig03_target_relevance_heatmap.png", "Target relevance (|Pearson| to sub-meter)"),
    ("fig04_flip_features_spotlight.png", "Flip-feature spotlight"),
    ("fig05_global_drop_partners.png", "Global drops vs greedy partner |Pearson|"),
    ("fig06_domain_survival.png", "Domain survival rates"),
    ("fig07_appliance_similarity.png", "Jaccard similarity between kept sets"),
    ("fig08_drop_decision_rules.png", "Drop decision rules (target vs priority)"),
]


def add_title_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle:
        tx = slide.placeholders[1].text_frame
        tx.text = subtitle


def add_bullets_slide(prs, title, bullets, max_lines=6):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        p.text = b
        p.level = 0
        p.font.size = Pt(14)


def add_image_slide(prs, title, image_path, caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    left = Inches(0.5)
    top = Inches(1.3)
    max_w = Inches(9)
    max_h = Inches(5.5)
    pic = slide.shapes.add_picture(str(image_path), left, top, width=max_w)
    if caption:
        tx = slide.shapes.add_textbox(Inches(0.5), top + pic.height + Inches(0.1), max_w, Inches(0.8))
        tf = tx.text_frame
        tf.text = caption
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT


def read_abstract():
    md = FS_DIR / "feature_selection.md"
    if not md.exists():
        return ""
    text = md.read_text(encoding="utf-8", errors="ignore")
    start = text.find("### Abstract")
    if start == -1:
        start = 0
    toc = text.find("## Table of contents", start)
    snippet = text[start:toc] if toc != -1 else text[start:start+1500]
    return "\n".join([l.strip() for l in snippet.splitlines() if l.strip()])[:2000]


def read_appendix_excerpt(lines=20):
    a = OUT_DIR / "stage01_summary.csv"
    if not a.exists():
        return ""
    # show top rows of stability tiers file if exists
    tiers = X_DIR / "feature_stability_tiers.csv"
    if tiers.exists():
        df = pd.read_csv(tiers) if 'pd' in sys.modules else None
    try:
        import pandas as pd
        df = pd.read_csv(tiers) if tiers.exists() else None
        if df is not None:
            return df.head(lines).to_string(index=False)
    except Exception:
        return ""


def add_appendix_slides(prs):
    # Add one slide pointing to appendix files + include the small table excerpt
    add_bullets_slide(prs, "Appendix — Audit files",
                      ["stage01_summary.csv (50×5 status pivot)",
                       "cross_appliance/global_drop_partners.csv",
                       "cross_appliance/feature_stability_tiers.csv",
                       "feature_selection/stage01_results_appendix.md"])
    excerpt = read_appendix_excerpt()
    if excerpt:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Appendix excerpt (stability tiers)"
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.0))
        tf = tx.text_frame
        for line in excerpt.splitlines():
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(10)


def main():
    prs = Presentation()
    add_title_slide(prs, "Stage 01 Feature Selection — Results (wk30)", "Auto-generated from feature_selection outputs")
    abstract = read_abstract()
    if abstract:
        add_bullets_slide(prs, "Abstract (excerpt)", textwrap.wrap(abstract, width=120)[:6])

    # Key summary slide
    add_bullets_slide(prs, "Key results (wk30)", [
        "Input HF features: 50",
        "Kept per appliance: 34 (16 dropped by correlation)",
        "Universal (kept 5/5): 30 features — recommended shared schema",
        "Global drops (0/5 kept): 10 features (structural redundancy)",
        "Flip features (appliance-dependent): 6 features — explained by target |r|"
    ])

    # Figures
    for fn, title in FIGS:
        img = X_DIR / fn
        capf = X_DIR / (fn.replace('.png', '_caption.txt'))
        caption = capf.read_text(encoding='utf-8') if capf.exists() else None
        if img.exists():
            add_image_slide(prs, title, img, caption=caption)
        else:
            add_bullets_slide(prs, f"{title} (missing)", [f"Expected image: {img}"])

    # Recommendations slide
    add_bullets_slide(prs, "Recommendations", [
        "Use Tier A (30 features) as shared HF schema for multi-appliance models.",
        "Use per-appliance 34-feature sets for best per-load performance.",
        "Run Stage 01 on multiple weeks (wk30+wk31) for stability on rare-ON loads.",
        "Proceed to Stage 02 (mRMR / multivariate) before final thesis features."
    ])

    # Appendix
    add_appendix_slides(prs)

    X_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(PPT_PATH)
    print('Wrote presentation: {}'.format(PPT_PATH))


if __name__ == '__main__':
    main()

